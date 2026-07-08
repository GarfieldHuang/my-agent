"""本地文件產生工具：Word（.docx）、PowerPoint（.pptx）、PDF。

以 function tool 掛進聊天迴圈，模型傳入結構化內容，
檔案直接寫到本機 generated_files/（agent 跑在本機，不需要下載）。
"""
from pathlib import Path

from .imagegen import _next_available

OUTPUT_DIR = Path("generated_files")

# docx / pdf 共用的內容區塊格式
_BLOCKS = {
    "type": "array",
    "description": "文件內容區塊，依序輸出",
    "items": {
        "type": "object",
        "properties": {
            "type":  {"type": "string",
                      "enum": ["heading", "paragraph", "bullets", "table"]},
            "text":  {"type": "string", "description": "heading / paragraph 的文字"},
            "level": {"type": "integer", "description": "heading 層級 1-4，預設 1"},
            "items": {"type": "array", "items": {"type": "string"},
                      "description": "bullets 的項目"},
            "rows":  {"type": "array",
                      "items": {"type": "array", "items": {"type": "string"}},
                      "description": "table 的列，第一列為表頭"},
        },
        "required": ["type"],
    },
}

_FILENAME = {"type": "string", "description": "輸出檔名，可不含副檔名"}

DOC_TOOLS = [
    {
        "type": "function",
        "name": "create_docx",
        "description": "產生 Word 文件（.docx）並存到使用者本機。適合報告、公文、會議記錄。",
        "parameters": {
            "type": "object",
            "properties": {
                "filename": _FILENAME,
                "title":    {"type": "string", "description": "文件標題（可省略）"},
                "blocks":   _BLOCKS,
            },
            "required": ["filename", "blocks"],
        },
    },
    {
        "type": "function",
        "name": "create_pdf",
        "description": "產生 PDF 文件並存到使用者本機。適合正式文件、可直接列印的版本。",
        "parameters": {
            "type": "object",
            "properties": {
                "filename": _FILENAME,
                "title":    {"type": "string", "description": "文件標題（可省略）"},
                "blocks":   _BLOCKS,
            },
            "required": ["filename", "blocks"],
        },
    },
    {
        "type": "function",
        "name": "create_pptx",
        "description": "產生 PowerPoint 簡報（.pptx）並存到使用者本機。",
        "parameters": {
            "type": "object",
            "properties": {
                "filename": _FILENAME,
                "title":    {"type": "string", "description": "封面標題（可省略）"},
                "subtitle": {"type": "string", "description": "封面副標題（可省略）"},
                "slides": {
                    "type": "array",
                    "description": "投影片列表",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title":   {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                            "notes":   {"type": "string",
                                        "description": "演講者備忘稿（可省略）"},
                        },
                        "required": ["title"],
                    },
                },
            },
            "required": ["filename", "slides"],
        },
    },
]


def is_doc_tool(name: str) -> bool:
    return name in _HANDLERS


def call_doc_tool(name: str, args: dict) -> tuple[str, Path | None]:
    """執行文件工具，回傳（給模型的結果訊息, 檔案路徑或 None）。"""
    try:
        path = _HANDLERS[name](args)
        return f"✓ 檔案已產生：{path}", path
    except Exception as e:
        return f"Error: {type(e).__name__}: {e}", None


def _out_path(filename: str, ext: str) -> Path:
    name = Path(str(filename)).name or f"document{ext}"   # 去除路徑，只留檔名
    if not name.lower().endswith(ext):
        name += ext
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return _next_available((OUTPUT_DIR / name).resolve())


# ── Word ─────────────────────────────────────────

def _create_docx(args: dict) -> Path:
    from docx import Document

    doc = Document()
    if args.get("title"):
        doc.add_heading(args["title"], 0)

    for b in args.get("blocks", []):
        t = b.get("type")
        if t == "heading":
            doc.add_heading(b.get("text", ""), min(max(b.get("level", 1), 1), 4))
        elif t == "paragraph":
            doc.add_paragraph(b.get("text", ""))
        elif t == "bullets":
            for item in b.get("items") or []:
                doc.add_paragraph(str(item), style="List Bullet")
        elif t == "table":
            rows = b.get("rows") or []
            if rows:
                table = doc.add_table(rows=len(rows),
                                      cols=max(len(r) for r in rows))
                table.style = "Table Grid"
                for i, r in enumerate(rows):
                    for j, cell in enumerate(r):
                        table.cell(i, j).text = str(cell)

    path = _out_path(args["filename"], ".docx")
    doc.save(str(path))
    return path


# ── PowerPoint ───────────────────────────────────

def _create_pptx(args: dict) -> Path:
    from pptx import Presentation

    prs = Presentation()
    if args.get("title"):
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = args["title"]
        if args.get("subtitle"):
            cover.placeholders[1].text = args["subtitle"]

    for sl in args.get("slides", []):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = sl.get("title", "")
        bullets = [str(b) for b in sl.get("bullets") or []]
        if bullets:
            tf = s.placeholders[1].text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                tf.add_paragraph().text = b
        if sl.get("notes"):
            s.notes_slide.notes_text_frame.text = str(sl["notes"])

    path = _out_path(args["filename"], ".pptx")
    prs.save(str(path))
    return path


# ── PDF ──────────────────────────────────────────

def _pdf_font() -> str:
    """挑一個支援中文的字型：微軟正黑體 → Acrobat 內建 CJK → Helvetica。"""
    from reportlab.pdfbase import pdfmetrics

    try:
        from reportlab.pdfbase.ttfonts import TTFont
        pdfmetrics.registerFont(
            TTFont("MSJhengHei", "C:/Windows/Fonts/msjh.ttc", subfontIndex=0))
        return "MSJhengHei"
    except Exception:
        pass
    try:
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        return "STSong-Light"
    except Exception:
        return "Helvetica"


def _create_pdf(args: dict) -> Path:
    from xml.sax.saxutils import escape

    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import (Paragraph, SimpleDocTemplate, Spacer,
                                    Table, TableStyle)

    font = _pdf_font()
    title_style = ParagraphStyle("t",  fontName=font, fontSize=20, leading=28,
                                 spaceAfter=14)
    body_style  = ParagraphStyle("b",  fontName=font, fontSize=11, leading=17,
                                 spaceAfter=6)
    h_sizes = {1: 16, 2: 14, 3: 12, 4: 12}

    story = []
    if args.get("title"):
        story.append(Paragraph(escape(args["title"]), title_style))

    for b in args.get("blocks", []):
        t = b.get("type")
        if t == "heading":
            size = h_sizes.get(min(max(b.get("level", 1), 1), 4), 14)
            story.append(Spacer(1, 6))
            story.append(Paragraph(
                escape(b.get("text", "")),
                ParagraphStyle("h", fontName=font, fontSize=size,
                               leading=size + 6, spaceAfter=8)))
        elif t == "paragraph":
            story.append(Paragraph(escape(b.get("text", "")), body_style))
        elif t == "bullets":
            for item in b.get("items") or []:
                story.append(Paragraph("• " + escape(str(item)), body_style))
        elif t == "table":
            rows = [[str(c) for c in r] for r in (b.get("rows") or [])]
            if rows:
                tbl = Table(rows, hAlign="LEFT")
                tbl.setStyle(TableStyle([
                    ("FONTNAME",   (0, 0), (-1, -1), font),
                    ("FONTSIZE",   (0, 0), (-1, -1), 10),
                    ("GRID",       (0, 0), (-1, -1), 0.5, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]))
                story.append(tbl)
                story.append(Spacer(1, 8))

    path = _out_path(args["filename"], ".pdf")
    SimpleDocTemplate(str(path), pagesize=A4).build(story)
    return path


_HANDLERS = {
    "create_docx": _create_docx,
    "create_pptx": _create_pptx,
    "create_pdf":  _create_pdf,
}
